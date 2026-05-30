#!/usr/bin/env python3
"""
ICTZoid Web Development Program — Session 01 Presentation
Generates a Google-Slides-compatible .pptx file.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUT = "/Users/samuel/Documents/Projects/lms/course-format/session-01-ictzoid.pptx"

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x08, 0x0C, 0x1A)
NAVY2     = RGBColor(0x0F, 0x16, 0x2A)
CYAN      = RGBColor(0x00, 0xD4, 0xFF)
WHITE     = RGBColor(0xF8, 0xFA, 0xFC)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)
DGRAY     = RGBColor(0x1E, 0x29, 0x3B)
LGRAY     = RGBColor(0x33, 0x41, 0x55)

# Slide size: 16:9 widescreen
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper utilities ──────────────────────────────────────────────────────────

def add_slide(layout_idx=6):
    """Blank slide (layout 6 = blank in most themes)."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, alpha=None):
    """Add a filled rectangle. l/t/w/h in Inches."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def txt(slide, text, l, t, w, h,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, font="Calibri"):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return txb

def txt_box(slide, lines, l, t, w, h,
            size=18, color=WHITE, line_space=1.15, font="Calibri"):
    """
    Add a multi-line text box.
    lines: list of (text, bold, color, size_override)
    """
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col, sz = item, False, color, size
        elif len(item) == 2:
            text, bold = item; col = color; sz = size
        elif len(item) == 3:
            text, bold, col = item; sz = size
        else:
            text, bold, col, sz = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.bold   = bold
        run.font.size   = Pt(sz)
        run.font.color.rgb = col
        run.font.name   = font
    return txb

def accent_bar(slide, l, t, w=0.06, h=7.3, color=CYAN):
    """Left vertical accent bar."""
    rect(slide, l, t, w, h, color)

def slide_label(slide, label, color=CYAN):
    """Small ALL-CAPS label top-right."""
    txt(slide, label,
        l=9.5, t=0.18, w=3.6, h=0.4,
        size=9, bold=True, color=color,
        align=PP_ALIGN.RIGHT, font="Calibri")

def divider_line(slide, l, t, w, color=DGRAY):
    rect(slide, l, t, w, 0.015, color)

def bullet_block(slide, items, l, t, w, h,
                 size=17, marker="▸", marker_color=CYAN, text_color=GRAY,
                 bold_first=False, font="Calibri"):
    """
    Render a list of bullet items.
    items: list of str OR (str, bold_flag) OR (str, bold_flag, color)
    """
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            label = item; bold = bold_first if first else False; col = text_color
        elif len(item) == 2:
            label, bold = item; col = text_color
        else:
            label, bold, col = item

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT

        # Marker run
        mr = p.add_run()
        mr.text = marker + "  "
        mr.font.bold  = True
        mr.font.size  = Pt(size)
        mr.font.color.rgb = marker_color
        mr.font.name  = font

        # Text run
        tr = p.add_run()
        tr.text = label
        tr.font.bold  = bold
        tr.font.size  = Pt(size)
        tr.font.color.rgb = col
        tr.font.name  = font
    return txb

def pill(slide, label, l, t, w=1.4, h=0.38,
         fill=None, border=CYAN, text_color=CYAN, size=13, font="Calibri"):
    """Rounded rectangle pill label."""
    if fill is None:
        fill = RGBColor(0x00, 0x1A, 0x22)
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.2)

    tf = shape.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = text_color
    run.font.name  = font
    return shape

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)

# Top cyan bar
rect(s, 0, 0, 13.33, 0.055, CYAN)

# Brand name
txt(s, "ICTZoid TechLab",
    l=0.35, t=0.55, w=8, h=0.55,
    size=14, bold=True, color=CYAN, font="Calibri")

# "WEB DEVELOPMENT PROGRAM" tag pill
pill(s, "WEB DEVELOPMENT PROGRAM",
     l=0.35, t=1.15, w=3.4, h=0.38,
     fill=RGBColor(0x00, 0x1E, 0x2A),
     border=CYAN, text_color=CYAN, size=11)

# Main headline
txt(s, "From Zero to",
    l=0.35, t=1.72, w=9, h=0.75,
    size=44, bold=True, color=WHITE, font="Calibri")
txt(s, "AI-Powered Developer",
    l=0.35, t=2.42, w=10, h=0.9,
    size=52, bold=True, color=CYAN, font="Calibri")

# Subline
txt(s, "6 Courses  ·  2 Master Sessions  ·  85+ Lessons  ·  Certificate",
    l=0.35, t=3.45, w=10, h=0.45,
    size=16, color=GRAY, font="Calibri")

# Divider
divider_line(s, 0.35, 4.0, 8.5, DGRAY)

# Author block
txt(s, "LEAD INSTRUCTOR",
    l=0.35, t=4.12, w=5, h=0.32,
    size=9, bold=True, color=GRAY, font="Calibri")
txt(s, "Samuel M. Ortil",
    l=0.35, t=4.42, w=5, h=0.55,
    size=24, bold=True, color=WHITE, font="Calibri")

# Course pills row
pills_data = [
    ("HTML",        CYAN,   RGBColor(0x00,0x1A,0x22)),
    ("CSS",         CYAN,   RGBColor(0x00,0x1A,0x22)),
    ("JavaScript",  CYAN,   RGBColor(0x00,0x1A,0x22)),
    ("Git",         CYAN,   RGBColor(0x00,0x1A,0x22)),
    ("React",       PURPLE, RGBColor(0x12,0x0A,0x2A)),
    ("Frontend AI", GREEN,  RGBColor(0x05,0x1A,0x0E)),
]
px = 0.35
for label, col, bg_col in pills_data:
    pw = max(1.1, len(label) * 0.125 + 0.6)
    pill(s, label, px, 5.22, pw, 0.38,
         fill=bg_col, border=col, text_color=col, size=13)
    px += pw + 0.12

# Year / cohort
txt(s, "May 2026  ·  First Session",
    l=0.35, t=6.8, w=6, h=0.38,
    size=11, color=GRAY, font="Calibri")

# Bottom bar
rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Who We Are
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "WHO WE ARE")

txt(s, "ICTZoid TechLab",
    l=0.35, t=0.32, w=9, h=0.7,
    size=38, bold=True, color=WHITE, font="Calibri")
txt(s, "A technology education lab built for the next generation of African developers",
    l=0.35, t=1.0, w=9.5, h=0.5,
    size=15, color=GRAY, font="Calibri")
divider_line(s, 0.35, 1.55, 12.6, DGRAY)

bullets_col1 = [
    ("Job-ready curriculum — everything maps to what real teams ship", False, WHITE),
    ("Cohort-based — you grow alongside a community, not in isolation", False, WHITE),
    ("Africa-first — built for our context, our pace, our ambitions", False, WHITE),
]
bullets_col2 = [
    ("Rigour — industry-standard tools and workflow", False, WHITE),
    ("Relevance — courses reflect what employers hire for today", False, WHITE),
    ("Community — learning is faster when surrounded by peers", False, WHITE),
]

bullet_block(s, bullets_col1,
             l=0.35, t=1.7, w=6.2, h=2.5,
             size=17, marker="▸", marker_color=CYAN)

# Pillars box
rect(s, 7.0, 1.6, 5.9, 3.1, NAVY2)
txt(s, "THREE PILLARS",
    l=7.2, t=1.7, w=5.5, h=0.38,
    size=9, bold=True, color=CYAN, font="Calibri")

py = 2.15
for pillar, col in [("Rigour", CYAN), ("Relevance", GREEN), ("Community", PURPLE)]:
    rect(s, 7.2, py, 5.5, 0.62, DGRAY)
    txt(s, pillar,
        l=7.38, t=py + 0.08, w=5.0, h=0.42,
        size=18, bold=True, color=col, font="Calibri")
    py += 0.78

txt(s, '"We do not teach theory for theory\'s sake."',
    l=0.35, t=4.7, w=12.5, h=0.52,
    size=16, italic=True, color=GRAY, align=PP_ALIGN.CENTER, font="Calibri")
rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Program
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "THE PROGRAM")

txt(s, "What You Are In",
    l=0.35, t=0.32, w=9, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
divider_line(s, 0.35, 0.98, 12.6, DGRAY)

# Stats row
stats = [("2", "Master Sessions"), ("6", "Courses"), ("85+", "Lessons"), ("1", "Certificate")]
sx = 0.5
for num, label in stats:
    rect(s, sx, 1.08, 2.7, 1.28, NAVY2)
    txt(s, num,   l=sx+0.12, t=1.15, w=2.4, h=0.72, size=42, bold=True, color=CYAN, font="Calibri")
    txt(s, label, l=sx+0.12, t=1.82, w=2.4, h=0.38, size=13, color=GRAY, font="Calibri")
    sx += 3.05

# Learning path
txt(s, "THE LEARNING PATH",
    l=0.35, t=2.55, w=7, h=0.35,
    size=9, bold=True, color=CYAN, font="Calibri")

path_items = [
    ("How the Internet Works", CYAN),
    ("SDLC",                   CYAN),
    ("HTML",                   WHITE),
    ("CSS",                    WHITE),
    ("JavaScript",             WHITE),
    ("Git & GitHub",           WHITE),
    ("React",                  PURPLE),
    ("Frontend AI",            GREEN),
]
px2 = 0.35
for i, (label, col) in enumerate(path_items):
    pw = max(1.1, len(label) * 0.13 + 0.55)
    pill(s, label, px2, 2.95, pw, 0.42,
         fill=DGRAY, border=col, text_color=col, size=13)
    px2 += pw + 0.06
    if i < len(path_items) - 1:
        txt(s, "→", l=px2 - 0.06, t=2.97, w=0.3, h=0.38,
            size=14, color=DGRAY, font="Calibri")
        px2 += 0.22

# Outcomes
txt(s, "WHERE YOU WILL BE AT THE END",
    l=0.35, t=3.62, w=8, h=0.35,
    size=9, bold=True, color=CYAN, font="Calibri")

outcomes = [
    "Building and deploying real web applications",
    "Integrating AI APIs (Claude, GPT-4o) inside production React apps",
    "Collaborating with version control the way professional teams do",
    "Holding a portfolio that proves it — not just a certificate",
]
bullet_block(s, outcomes,
             l=0.35, t=3.98, w=12.5, h=2.5,
             size=17, marker="✓", marker_color=GREEN)
rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Why Trust Me
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "YOUR INSTRUCTOR")

txt(s, "Samuel M. Ortil",
    l=0.35, t=0.28, w=9, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
txt(s, "Lead Instructor — ICTZoid TechLab",
    l=0.35, t=0.92, w=8, h=0.4,
    size=15, color=GRAY, font="Calibri")
divider_line(s, 0.35, 1.38, 12.6, DGRAY)

# Teaching column
txt(s, "TEACHING & MENTORING",
    l=0.35, t=1.5, w=6, h=0.32,
    size=9, bold=True, color=CYAN, font="Calibri")
teaching = [
    ("2019 — Web Design · Classic InfoTech, FUTA Resource Centre", False, WHITE),
    ("2019 — Web Design Instructor · Obafemi Awolowo University (OAU)", False, WHITE),
    ("Mentor — Top Universe", False, WHITE),
    ("Mentor — Maxify Global", False, WHITE),
]
bullet_block(s, teaching,
             l=0.35, t=1.82, w=6.1, h=2.1,
             size=15, marker="▸", marker_color=CYAN)

# International experience column
txt(s, "INTERNATIONAL EXPERIENCE",
    l=6.8, t=1.5, w=6, h=0.32,
    size=9, bold=True, color=GREEN, font="Calibri")

roles = [
    ("CTO",              "Maxify Global",  "Nigeria · USA · Canada", CYAN),
    ("Senior Engineer",  "Breezelearn",    "Canada",                 GREEN),
    ("Fullstack Dev",    "CircleHQ",       "Germany",                PURPLE),
]
ry = 1.85
for role, company, location, col in roles:
    rect(s, 6.8, ry, 6.15, 0.78, DGRAY)
    txt(s, role,     l=7.0, t=ry+0.04, w=5.8, h=0.32, size=10, bold=True, color=GRAY,  font="Calibri")
    txt(s, company,  l=7.0, t=ry+0.24, w=5.8, h=0.32, size=17, bold=True, color=col,   font="Calibri")
    txt(s, location, l=7.0, t=ry+0.50, w=5.8, h=0.26, size=12, color=GRAY, font="Calibri")
    ry += 0.88

# What this means box
rect(s, 0.35, 4.55, 12.55, 1.38, NAVY2)
txt(s, "WHAT THIS MEANS FOR YOU",
    l=0.6, t=4.65, w=8, h=0.32,
    size=9, bold=True, color=CYAN, font="Calibri")
means = [
    "I have built what I'm teaching — not read about it",
    "I have taught beginners before — I know where people get stuck",
    "I have worked in companies across 4 countries — I know what employers expect",
]
bullet_block(s, means,
             l=0.55, t=4.95, w=12.1, h=0.9,
             size=14, marker="→", marker_color=GREEN)
rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Requirements
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "REQUIREMENTS")

txt(s, "What You Need to Begin",
    l=0.35, t=0.32, w=9, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
divider_line(s, 0.35, 0.98, 12.6, DGRAY)

reqs = [
    ("01", "GitHub Account",
     "Go to github.com · Create a free account · Use a professional username\nThis becomes your developer identity and portfolio home",
     CYAN),
    ("02", "A Device",
     "Laptop or desktop — Windows, macOS, or Linux\nStable internet connection for live sessions",
     GREEN),
    ("03", "VS Code (Text Editor)",
     "Download free at code.visualstudio.com\nThis is the industry-standard editor we will use throughout",
     PURPLE),
    ("04", "A Learning Mindset",
     "No prior coding experience required\nConsistency matters more than talent — show up every session",
     ORANGE),
]

rx, ry = 0.35, 1.1
for i, (num, title, body, col) in enumerate(reqs):
    col_x = rx + (i % 2) * 6.5
    row_y = ry + (i // 2) * 2.55
    rect(s, col_x, row_y, 6.1, 2.3, NAVY2)
    txt(s, num,   l=col_x+0.18, t=row_y+0.12, w=1,   h=0.5, size=28, bold=True, color=col, font="Calibri")
    txt(s, title, l=col_x+0.18, t=row_y+0.6,  w=5.7, h=0.45, size=18, bold=True, color=WHITE, font="Calibri")
    txt(s, body,  l=col_x+0.18, t=row_y+1.04, w=5.7, h=1.05, size=13, color=GRAY, font="Calibri")

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Curriculum (Part 1 — Master Sessions)
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "CURRICULUM — MASTER SESSIONS")

txt(s, "Course Curriculum",
    l=0.35, t=0.28, w=9, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
txt(s, "8 units in total  ·  Part 1 of 2",
    l=0.35, t=0.92, w=9, h=0.38,
    size=14, color=GRAY, font="Calibri")
divider_line(s, 0.35, 1.3, 12.6, DGRAY)

# Master session cards
ms_cards = [
    ("Master Session 1", "How the Internet Works",
     ["DNS, TCP/IP, HTTP, packets, status codes", "URL anatomy, TLS / HTTPS", "How browsers render pages"],
     CYAN, "No prerequisite"),
    ("Master Session 2", "Software Development Lifecycle",
     ["Agile, Scrum, sprint planning, retrospectives", "CI/CD pipelines, DevOps principles", "Team roles, burn-down charts, velocity"],
     CYAN, "No prerequisite"),
]
mx = 0.35
for tag, title, points, col, prereq in ms_cards:
    rect(s, mx, 1.45, 6.1, 4.95, NAVY2)
    pill(s, tag, mx+0.2, 1.58, 2.4, 0.35,
         fill=DGRAY, border=col, text_color=col, size=11)
    txt(s, title,
        l=mx+0.2, t=2.02, w=5.65, h=0.75,
        size=20, bold=True, color=WHITE, font="Calibri")
    divider_line(s, mx+0.2, 2.78, 5.5, DGRAY)
    by = 2.92
    for pt in points:
        txt(s, "▸  " + pt,
            l=mx+0.2, t=by, w=5.65, h=0.4,
            size=14, color=GRAY, font="Calibri")
        by += 0.45
    txt(s, "⚙  " + prereq,
        l=mx+0.2, t=5.98, w=5.65, h=0.35,
        size=12, color=GREEN, font="Calibri")
    mx += 6.55

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Curriculum (Part 2 — Courses)
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "CURRICULUM — COURSES")

txt(s, "6 Courses — The Full Path",
    l=0.35, t=0.28, w=9, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
divider_line(s, 0.35, 0.92, 12.6, DGRAY)

courses = [
    ("01", "HTML",            "Semantic structure, accessibility, ARIA, forms",    CYAN),
    ("02", "CSS",             "Flexbox, Grid, responsive design, animations",       CYAN),
    ("03", "JavaScript",      "DOM, events, async/await, Fetch API, localStorage",  CYAN),
    ("04", "Git & GitHub",    "Commits, branches, pull requests, CI pipelines",     CYAN),
    ("05", "React",           "Hooks, routing, state, TanStack Query, Tailwind",    PURPLE),
    ("06", "Frontend AI",     "Claude & GPT-4o APIs, streaming, voice, tool call",  GREEN),
]

cx, cy = 0.35, 1.0
for i, (num, title, desc, col) in enumerate(courses):
    row = i // 3
    col_i = i % 3
    bx = cx + col_i * 4.32
    by = cy + row * 2.75
    rect(s, bx, by, 4.1, 2.5, NAVY2)
    txt(s, num,   l=bx+0.18, t=by+0.14, w=0.8, h=0.5, size=22, bold=True, color=col, font="Calibri")
    txt(s, title, l=bx+0.18, t=by+0.62, w=3.7, h=0.42, size=17, bold=True, color=WHITE, font="Calibri")
    divider_line(s, bx+0.18, by+1.1, 3.55, DGRAY)
    txt(s, desc,  l=bx+0.18, t=by+1.22, w=3.7, h=1.05, size=12, color=GRAY, font="Calibri")

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Mode of Learning
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "MODE OF LEARNING")

txt(s, "How This Program Is Delivered",
    l=0.35, t=0.28, w=10, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
txt(s, '"This is not a YouTube playlist. This is a structured cohort."',
    l=0.35, t=0.92, w=12, h=0.42,
    size=14, italic=True, color=GRAY, font="Calibri")
divider_line(s, 0.35, 1.38, 12.6, DGRAY)

modes = [
    ("📚", "LMS Platform",          "All lessons, notes, exercises & progress tracking in one place",  CYAN),
    ("📖", "Program Ebook",          "Complete written reference — read offline, print, annotate",      CYAN),
    ("🐙", "GitHub Notes",           "Version-controlled notes — always up to date, forkable",          CYAN),
    ("💬", "Learning Community",     "WhatsApp group — questions, accountability, peer support",         GREEN),
    ("🎥", "Weekend Live Sessions",  "Real-time coding, Q&A, recorded for replay in the LMS",           GREEN),
    ("🛠", "Workshops",              "Hands-on project sessions — build real things with live feedback",  GREEN),
    ("🎓", "Master Sessions",        "Full-cohort deep dives — How the Web Works + SDLC",               PURPLE),
]

mx2, my2 = 0.35, 1.52
for i, (icon, title, desc, col) in enumerate(modes):
    row2  = i // 4
    col_i = i % 4
    bx2   = mx2 + col_i * 3.24
    by2   = my2 + row2 * 2.75
    rect(s, bx2, by2, 3.04, 2.5, NAVY2)
    txt(s, icon,  l=bx2+0.15, t=by2+0.15, w=0.7,  h=0.45, size=22, color=col,   font="Segoe UI Emoji")
    txt(s, title, l=bx2+0.15, t=by2+0.62, w=2.72, h=0.42, size=14, bold=True, color=WHITE, font="Calibri")
    divider_line(s, bx2+0.15, by2+1.08, 2.6, DGRAY)
    txt(s, desc,  l=bx2+0.15, t=by2+1.2,  w=2.72, h=1.12, size=11, color=GRAY,  font="Calibri")

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Program Fee
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "PROGRAM FEE")

txt(s, "Program Fee",
    l=0.35, t=0.28, w=9, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
divider_line(s, 0.35, 0.92, 12.6, DGRAY)

# Free banner
rect(s, 0.35, 1.05, 12.55, 1.55, RGBColor(0x00, 0x1A, 0x12))
txt(s, "FREE",
    l=0.6, t=1.12, w=3, h=1.05,
    size=62, bold=True, color=GREEN, font="Calibri")
txt(s, "The program itself costs nothing.",
    l=3.5, t=1.2, w=9, h=0.45,
    size=18, bold=True, color=WHITE, font="Calibri")
txt(s, "We believe access to quality technical education should not depend on your bank balance.",
    l=3.5, t=1.62, w=9, h=0.4,
    size=13, color=GRAY, font="Calibri")

# What is paid
txt(s, "WHAT IS PAID FOR",
    l=0.35, t=2.78, w=6, h=0.32,
    size=9, bold=True, color=ORANGE, font="Calibri")

paid_items = [
    ("Ebook (full program)", "Production cost for the authored & designed written guide"),
    ("Certificate",          "Verified, designed, and archived certificate of completion"),
    ("Workshop Materials",   "Printed / digital materials for intensive workshops"),
]
py2 = 3.12
for item, detail in paid_items:
    rect(s, 0.35, py2, 12.55, 0.75, DGRAY)
    txt(s, "💰  " + item,  l=0.55, t=py2+0.08, w=5,    h=0.35, size=15, bold=True, color=WHITE,  font="Calibri")
    txt(s, detail,          l=6.0,  t=py2+0.08, w=6.65, h=0.55, size=13, color=GRAY, font="Calibri")
    py2 += 0.85

txt(s, "No hidden fees. No surprise charges. Costs are communicated clearly before any purchase.",
    l=0.35, t=5.72, w=12.5, h=0.42,
    size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER, font="Calibri")

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Scholarship
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "SCHOLARSHIP")

txt(s, "Scholarship Programme",
    l=0.35, t=0.28, w=10, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
txt(s, '"No one should be left out because of finances."',
    l=0.35, t=0.92, w=12, h=0.42,
    size=14, italic=True, color=GRAY, font="Calibri")
divider_line(s, 0.35, 1.38, 12.6, DGRAY)

# Steps
steps = [
    ("01", "Fill the Application Form",
     "Scholarship Application Form: [link to be inserted]\nTakes approximately 5 minutes to complete",
     CYAN),
    ("02", "What the Form Asks",
     "Your name & contact details\nWhich resources you need support with\n3–5 sentence statement on your commitment",
     CYAN),
    ("03", "Wait for a Response",
     "Applications reviewed within 5 business days\nNotified privately via email or WhatsApp\nDecisions made with dignity and privacy",
     GREEN),
]
sx2 = 0.35
for num, title, body, col in steps:
    rect(s, sx2, 1.55, 4.1, 4.0, NAVY2)
    txt(s, num,   l=sx2+0.2, t=1.68, w=1, h=0.65, size=32, bold=True, color=col, font="Calibri")
    txt(s, title, l=sx2+0.2, t=2.3,  w=3.7, h=0.55, size=16, bold=True, color=WHITE, font="Calibri")
    divider_line(s, sx2+0.2, 2.9, 3.55, DGRAY)
    txt(s, body,  l=sx2+0.2, t=3.0,  w=3.7, h=2.3, size=13, color=GRAY, font="Calibri")
    sx2 += 4.35

# Note
rect(s, 0.35, 5.7, 12.55, 0.75, RGBColor(0x00, 0x14, 0x1A))
txt(s, "Do not wait for the scholarship decision to start. The program content is FREE — begin immediately.",
    l=0.55, t=5.78, w=12.1, h=0.55,
    size=14, italic=True, color=CYAN, align=PP_ALIGN.CENTER, font="Calibri")

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — What Happens Next
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
accent_bar(s, 0, 0, 0.07, 7.5, CYAN)
rect(s, 0, 0, 13.33, 0.055, CYAN)
slide_label(s, "WHAT HAPPENS NEXT")

txt(s, "Starting from Today",
    l=0.35, t=0.28, w=9, h=0.65,
    size=36, bold=True, color=WHITE, font="Calibri")
divider_line(s, 0.35, 0.92, 12.6, DGRAY)

timeline = [
    ("Today",           "First session — orientation, introductions, setup",              CYAN),
    ("This Weekend",    "Master Session 1 — How the Internet Works",                     CYAN),
    ("Next Weekend",    "Master Session 2 — Software Development Lifecycle",             CYAN),
    ("Following Weeks", "HTML → CSS → JavaScript → Git → React → Frontend AI",          GREEN),
    ("End of Program",  "Portfolio review, certificate, next steps",                     PURPLE),
]
ty2 = 1.08
for period, desc, col in timeline:
    rect(s, 0.35, ty2, 12.55, 0.78, NAVY2)
    rect(s, 0.35, ty2, 0.08, 0.78, col)
    txt(s, period, l=0.6,  t=ty2+0.05, w=2.8, h=0.65, size=14, bold=True, color=col,   font="Calibri")
    txt(s, desc,   l=3.55, t=ty2+0.16, w=9.2, h=0.45, size=15, color=WHITE, font="Calibri")
    ty2 += 0.9

# Action items
txt(s, "YOUR ACTION ITEMS BEFORE NEXT SESSION",
    l=0.35, t=5.65, w=9, h=0.32,
    size=9, bold=True, color=ORANGE, font="Calibri")
actions = [
    "Create / confirm your GitHub account",
    "Download and install VS Code",
    "Introduce yourself in the WhatsApp group",
    "Read the program overview on the LMS",
]
ax = 0.35
for act in actions:
    pill(s, "☐  " + act, ax, 5.98, 3.0, 0.38,
         fill=DGRAY, border=ORANGE, text_color=WHITE, size=12)
    ax += 3.25

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Q&A / Closing
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, 13.33, 0.055, CYAN)

# Large Q&A
txt(s, "Q&A",
    l=0, t=1.3, w=13.33, h=2.5,
    size=120, bold=True, color=RGBColor(0x0F,0x1E,0x38),
    align=PP_ALIGN.CENTER, font="Calibri")

txt(s, "Questions?",
    l=0, t=3.2, w=13.33, h=0.75,
    size=36, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER, font="Calibri")

txt(s, '"You don\'t have to be great to start. But you have to start to be great."',
    l=1.5, t=4.1, w=10.33, h=0.6,
    size=16, italic=True, color=GRAY,
    align=PP_ALIGN.CENTER, font="Calibri")

divider_line(s, 3.0, 4.88, 7.33, DGRAY)

channels = [
    ("LMS", "All lessons, notes & progress"),
    ("WhatsApp", "Community & quick questions"),
    ("GitHub", "Notes & code — always latest"),
    ("Ebook", "Full written reference"),
]
cx2 = 1.5
for ch, desc in channels:
    txt(s, ch,   l=cx2, t=5.05, w=2.5, h=0.38, size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER, font="Calibri")
    txt(s, desc, l=cx2, t=5.42, w=2.5, h=0.38, size=11, color=GRAY, align=PP_ALIGN.CENTER, font="Calibri")
    cx2 += 2.6

txt(s, "ICTZoid TechLab  ·  Web Development Program  ·  2026  ·  Samuel M. Ortil",
    l=0, t=6.85, w=13.33, h=0.38,
    size=11, color=GRAY, align=PP_ALIGN.CENTER, font="Calibri")

rect(s, 0, 7.38, 13.33, 0.12, CYAN)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
